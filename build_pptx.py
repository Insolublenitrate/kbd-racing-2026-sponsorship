from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(18, 18, 18)

prs = Presentation()
# Set 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]

# SLIDE 1: Cover
slide1 = prs.slides.add_slide(blank_slide_layout)
add_dark_bg(slide1)
slide1.shapes.add_picture('assets/FB_IMG_1780578610600.jpg', 0, 0, width=prs.slide_width, height=prs.slide_height)
# Dark overlay
overlay = slide1.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(18, 18, 18)
overlay.fill.transparency = 0.5
overlay.line.fill.background()

txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1))
tf = txBox.text_frame
tf.text = "KIND OF A BIG DILL RACING"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.size = Pt(64)
tf.paragraphs[0].font.name = 'Impact'
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

txBox2 = slide1.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
tf2 = txBox2.text_frame
tf2.text = "2026 SPONSORSHIP PROPOSAL"
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
tf2.paragraphs[0].font.size = Pt(48)
tf2.paragraphs[0].font.name = 'Impact'
tf2.paragraphs[0].font.color.rgb = RGBColor(57, 255, 20)

txBox3 = slide1.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(1))
tf3 = txBox3.text_frame
tf3.text = "Precision Engineering Meets High-Performance Speed"
tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
tf3.paragraphs[0].font.size = Pt(24)
tf3.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# SLIDE 2: About the Team
slide2 = prs.slides.add_slide(blank_slide_layout)
add_dark_bg(slide2)
slide2.shapes.add_picture('assets/FB_IMG_1780578645493.jpg', 0, 0, height=prs.slide_height)

def add_heading(slide, text, left, top):
    tx = slide.shapes.add_textbox(left, top, Inches(8), Inches(1))
    tf = tx.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.name = 'Impact'
    tf.paragraphs[0].font.color.rgb = RGBColor(57, 255, 20)

add_heading(slide2, "ABOUT THE TEAM", Inches(5), Inches(0.5))

tx = slide2.shapes.add_textbox(Inches(5), Inches(1.5), Inches(8), Inches(5))
tf = tx.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Team Principal: Dillon Smith"
p.font.size = Pt(32)
p.font.name = 'Impact'
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "U.S. Air Force Veteran (2W251 Nuclear Weapons Specialist)\nSenior Product Engineer"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "\nThe Riders:"
p.font.size = Pt(28)
p.font.name = 'Impact'
p.font.color.rgb = RGBColor(255, 255, 255)

for rider in ["Dillon 'Big Dill' Smith (Lead Rider)", "Dylan 'Kentucky Kid' Yelton (Amateur Rider)", "Andrew 'Hooligan' Perry (Amateur Rider)"]:
    p = tf.add_paragraph()
    p.text = "• " + rider
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "\nThe Machines:"
p.font.size = Pt(28)
p.font.name = 'Impact'
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "Premium 2021 Ducati Panigale V2 (955cc twin)"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(255, 255, 255)

# SLIDE 3: Engineering Meets Telemetry
slide3 = prs.slides.add_slide(blank_slide_layout)
add_dark_bg(slide3)
add_heading(slide3, "ENGINEERING MEETS TELEMETRY", Inches(0.5), Inches(0.5))

tx = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7), Inches(5))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The Data-Driven Competitive Edge"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(255, 255, 255)

bullets = [
    "Applying rigorous product engineering principles directly to racetrack performance.",
    "Consistently running an elite, competitive 2:00 - 2:02 amateur pace at NCM Motorsports Park.",
    "Tracking precise lap-telemetry using custom-integrated RaceBox GPS hardware.",
    "Processing raw data through our custom-developed 'KBD Track Engineer App' to optimize lines."
]
for b in bullets:
    p = tf.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)

slide3.shapes.add_picture('assets/Screenshot 1.jpg', Inches(8), Inches(1.5), height=Inches(4.5))
slide3.shapes.add_picture('assets/Screenshot 2.jpg', Inches(10.5), Inches(1.5), height=Inches(4.5))

# SLIDE 4: Podium Performance
slide4 = prs.slides.add_slide(blank_slide_layout)
add_dark_bg(slide4)
slide4.shapes.add_picture('assets/FB_IMG_1780578614560.jpg', 0, 0, height=prs.slide_height)
slide4.shapes.add_picture('assets/_B0A4077.jpg', Inches(2.5), 0, height=prs.slide_height)

add_heading(slide4, "PODIUM PERFORMANCE & IMPACT", Inches(5.5), Inches(0.5))

tx = slide4.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(7.5), Inches(5))
tf = tx.text_frame
tf.word_wrap = True

bullets = [
    "Over 15 amateur podium finishes across three competitive race organizations.",
    "Growing brand presence with thousands of active social media supporters.",
    "Highly visible paddock presence at N2 Track Days, Precision Trackdays, and Fast Line Track Days.",
    "Dedicated off-season indoor kart track training and racing aboard a 50CC CRF mini-bike."
]
for b in bullets:
    p = tf.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "\nCurrent Partners:"
p.font.size = Pt(28)
p.font.name = 'Impact'
p.font.color.rgb = RGBColor(255, 255, 255)

p = tf.add_paragraph()
p.text = "Official Technical Sponsor: RaceBox Pro (Promo Code: KBD15)"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(255, 255, 255)

# SLIDE 5: Partner With Us
slide5 = prs.slides.add_slide(blank_slide_layout)
add_dark_bg(slide5)
slide5.shapes.add_picture('assets/FB_IMG_1780578645493.jpg', 0, 0, width=prs.slide_width, height=prs.slide_height)
overlay = slide5.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(18, 18, 18)
overlay.fill.transparency = 0.5
overlay.line.fill.background()

tx = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(6))
tf = tx.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "PARTNER WITH US"
p.font.size = Pt(48)
p.font.name = 'Impact'
p.font.color.rgb = RGBColor(57, 255, 20)
p.alignment = PP_ALIGN.CENTER

def add_tier(title, desc):
    p = tf.add_paragraph()
    p.text = "\n" + title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER

add_tier("Title Sponsor", "Exclusive motorcycle livery branding, premier apparel placement, dedicated social media campaigns, and title paddock-banner display.")
add_tier("Associate Sponsor", "Secondary livery branding, apparel logos, and consistent social media mentions.")
add_tier("Product Sponsor", "Paddock equipment/apparel showcases, paddock-social-hour product sampling, and logo placement.")

p = tf.add_paragraph()
p.text = "\nCONTACT US\nDillon Smith | insolublenitrate@gmail.com | (513) 926-1478"
p.font.size = Pt(24)
p.font.name = 'Impact'
p.font.color.rgb = RGBColor(57, 255, 20)
p.alignment = PP_ALIGN.CENTER

prs.save('KBD_Racing_2026_Sponsorship_Deck.pptx')
print("Successfully generated KBD_Racing_2026_Sponsorship_Deck.pptx")
